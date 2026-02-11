"""
=================================================================
AI POWERPOINT GENERATOR PRO - COMPLETE SYSTEM
Login System + Admin Panel (Create/Edit/Delete Users) + Full Dashboards
=================================================================
"""

import streamlit as st
import requests
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import time
import json
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from datetime import datetime
import hashlib
import sqlite3
import plotly.express as px
import plotly.graph_objects as go

# ============================================================================
# DATABASE INITIALIZATION & USER MANAGEMENT
# ============================================================================

def init_database():
    """Initialize SQLite database with all tables"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    
    # Users table with all fields
    c.execute('''CREATE TABLE IF NOT EXISTS users
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  username TEXT UNIQUE NOT NULL,
                  password_hash TEXT NOT NULL,
                  email TEXT,
                  full_name TEXT,
                  phone TEXT,
                  created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  last_login TIMESTAMP,
                  is_active BOOLEAN DEFAULT 1,
                  role TEXT DEFAULT 'user')''')
    
    # Usage logs table
    c.execute('''CREATE TABLE IF NOT EXISTS usage_logs
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  action TEXT,
                  topic TEXT,
                  slides_count INTEGER,
                  timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    # Sessions table
    c.execute('''CREATE TABLE IF NOT EXISTS sessions
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  login_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  logout_time TIMESTAMP,
                  is_active BOOLEAN DEFAULT 1,
                  session_token TEXT,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    # Templates table
    c.execute('''CREATE TABLE IF NOT EXISTS templates
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  template_name TEXT,
                  template_data TEXT,
                  created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    # Migration: Add missing columns
    try:
        c.execute("PRAGMA table_info(sessions)")
        columns = [column[1] for column in c.fetchall()]
        if 'is_active' not in columns:
            c.execute("ALTER TABLE sessions ADD COLUMN is_active BOOLEAN DEFAULT 1")
        if 'session_token' not in columns:
            c.execute("ALTER TABLE sessions ADD COLUMN session_token TEXT")
        conn.commit()
    except:
        pass
    
    # Create default admin user
    c.execute("SELECT * FROM users WHERE username = 'admin'")
    if not c.fetchone():
        admin_password = hashlib.sha256('admin123'.encode()).hexdigest()
        c.execute("""INSERT INTO users (username, password_hash, email, full_name, role) 
                     VALUES (?, ?, ?, ?, ?)""",
                  ('admin', admin_password, 'admin@pptgen.com', 'System Administrator', 'admin'))
    
    conn.commit()
    conn.close()

def hash_password(password):
    """Hash password using SHA256"""
    return hashlib.sha256(password.encode()).hexdigest()

def verify_user(username, password):
    """Verify user credentials and create session"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    
    try:
        password_hash = hash_password(password)
        c.execute("""SELECT id, username, role, is_active, email, full_name 
                     FROM users WHERE username = ? AND password_hash = ?""",
                  (username, password_hash))
        
        user = c.fetchone()
        
        if user and user[3]:  # Check if user exists and is active
            user_id = user[0]
            c.execute("UPDATE users SET last_login = ? WHERE id = ?", (datetime.now(), user_id))
            session_token = hashlib.md5(f"{user_id}{datetime.now()}".encode()).hexdigest()
            
            try:
                c.execute("""UPDATE sessions SET is_active = 0, logout_time = ? 
                             WHERE user_id = ? AND is_active = 1""", 
                          (datetime.now(), user_id))
            except:
                pass
            
            try:
                c.execute("""INSERT INTO sessions (user_id, login_time, is_active, session_token) 
                             VALUES (?, ?, ?, ?)""",
                          (user_id, datetime.now(), 1, session_token))
            except:
                c.execute("INSERT INTO sessions (user_id, login_time) VALUES (?, ?)",
                          (user_id, datetime.now()))
            
            conn.commit()
            conn.close()
            
            return {
                'id': user[0],
                'username': user[1],
                'role': user[2],
                'is_active': user[3],
                'email': user[4] or '',
                'full_name': user[5] or username,
                'session_token': session_token
            }
        
        conn.close()
        return None
    except Exception as e:
        conn.close()
        return None

def create_user(username, password, email, full_name, phone, role='user'):
    """Create new user (admin only)"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    try:
        password_hash = hash_password(password)
        c.execute("""INSERT INTO users (username, password_hash, email, full_name, phone, role) 
                     VALUES (?, ?, ?, ?, ?, ?)""",
                  (username, password_hash, email, full_name, phone, role))
        conn.commit()
        conn.close()
        return True, "User created successfully"
    except sqlite3.IntegrityError:
        conn.close()
        return False, "Username already exists"
    except Exception as e:
        conn.close()
        return False, str(e)

def update_user(user_id, username, email, full_name, phone, is_active, password=None):
    """Update existing user"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    try:
        if password:
            password_hash = hash_password(password)
            c.execute("""UPDATE users 
                         SET username = ?, email = ?, full_name = ?, phone = ?, 
                             is_active = ?, password_hash = ?
                         WHERE id = ?""",
                      (username, email, full_name, phone, is_active, password_hash, user_id))
        else:
            c.execute("""UPDATE users 
                         SET username = ?, email = ?, full_name = ?, phone = ?, is_active = ?
                         WHERE id = ?""",
                      (username, email, full_name, phone, is_active, user_id))
        conn.commit()
        conn.close()
        return True, "User updated successfully"
    except sqlite3.IntegrityError:
        conn.close()
        return False, "Username already exists"
    except Exception as e:
        conn.close()
        return False, str(e)

def delete_user(user_id):
    """Delete user"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    try:
        c.execute("DELETE FROM users WHERE id = ?", (user_id,))
        conn.commit()
        conn.close()
        return True, "User deleted successfully"
    except Exception as e:
        conn.close()
        return False, str(e)

def get_user_by_id(user_id):
    """Get user details by ID"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""SELECT id, username, email, full_name, phone, is_active, role, created_at, last_login
                 FROM users WHERE id = ?""", (user_id,))
    user = c.fetchone()
    conn.close()
    return user

def get_all_users():
    """Get all users"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""SELECT id, username, email, full_name, phone, created_at, last_login, is_active, role 
                 FROM users ORDER BY created_at DESC""")
    users = c.fetchall()
    conn.close()
    return users

def get_currently_logged_in_users():
    """Get currently logged in users"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    try:
        c.execute("""
            SELECT u.id, u.username, u.email, u.full_name, s.login_time, u.role
            FROM sessions s
            JOIN users u ON s.user_id = u.id
            WHERE s.is_active = 1
            ORDER BY s.login_time DESC
        """)
        active_users = c.fetchall()
    except:
        active_users = []
    conn.close()
    return active_users

def logout_user(user_id):
    """Logout user"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    try:
        c.execute("""UPDATE sessions SET is_active = 0, logout_time = ? 
                     WHERE user_id = ? AND is_active = 1""",
                  (datetime.now(), user_id))
    except:
        c.execute("UPDATE sessions SET logout_time = ? WHERE user_id = ? AND logout_time IS NULL",
                  (datetime.now(), user_id))
    conn.commit()
    conn.close()

def log_usage(user_id, action, topic="", slides_count=0):
    """Log user activity"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""INSERT INTO usage_logs (user_id, action, topic, slides_count) 
                 VALUES (?, ?, ?, ?)""",
              (user_id, action, topic, slides_count))
    conn.commit()
    conn.close()

def get_user_stats(user_id):
    """Get user statistics"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM usage_logs WHERE user_id = ? AND action = 'generate_presentation'", (user_id,))
    total_presentations = c.fetchone()[0]
    c.execute("SELECT SUM(slides_count) FROM usage_logs WHERE user_id = ? AND action = 'generate_presentation'", (user_id,))
    total_slides = c.fetchone()[0] or 0
    c.execute("SELECT COUNT(*) FROM sessions WHERE user_id = ?", (user_id,))
    total_logins = c.fetchone()[0]
    conn.close()
    return {'total_presentations': total_presentations, 'total_slides': total_slides, 'total_logins': total_logins}

def get_user_activity_details(user_id):
    """Get detailed activity for a user"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""SELECT action, topic, slides_count, timestamp
                 FROM usage_logs WHERE user_id = ?
                 ORDER BY timestamp DESC LIMIT 20""", (user_id,))
    activities = c.fetchall()
    conn.close()
    return activities

def get_all_user_activities():
    """Get all system activities"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""SELECT u.username, l.action, l.topic, l.slides_count, l.timestamp
                 FROM usage_logs l
                 JOIN users u ON l.user_id = u.id
                 ORDER BY l.timestamp DESC LIMIT 100""")
    activities = c.fetchall()
    conn.close()
    return activities

def get_system_stats():
    """Get system-wide statistics"""
    conn = sqlite3.connect('ppt_system.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM users WHERE role = 'user'")
    total_users = c.fetchone()[0]
    try:
        c.execute("SELECT COUNT(*) FROM sessions WHERE is_active = 1")
        currently_online = c.fetchone()[0]
    except:
        currently_online = 0
    c.execute("SELECT COUNT(*) FROM usage_logs WHERE action = 'generate_presentation'")
    total_presentations = c.fetchone()[0]
    c.execute("SELECT SUM(slides_count) FROM usage_logs WHERE action = 'generate_presentation'")
    total_slides = c.fetchone()[0] or 0
    c.execute("SELECT COUNT(*) FROM sessions WHERE DATE(login_time) = DATE('now')")
    today_logins = c.fetchone()[0]
    conn.close()
    return {
        'total_users': total_users,
        'currently_online': currently_online,
        'total_presentations': total_presentations,
        'total_slides': total_slides,
        'today_logins': today_logins
    }

# ============================================================================
# TEMPLATE MANAGEMENT
# ============================================================================

def generate_template_id():
    return hashlib.md5(str(datetime.now().timestamp()).encode()).hexdigest()[:8]

def save_template_to_state(name, template_data):
    template_id = generate_template_id()
    template_data['id'] = template_id
    template_data['name'] = name
    template_data['created_at'] = datetime.now().strftime("%Y-%m-%d %H:%M")
    template_data['usage_count'] = 0
    st.session_state.templates[template_id] = template_data
    return template_id

def delete_template(template_id):
    if template_id in st.session_state.templates:
        del st.session_state.templates[template_id]
        return True
    return False

def export_all_templates():
    return json.dumps(st.session_state.templates, indent=2)

def import_templates(json_data):
    try:
        templates = json.loads(json_data)
        st.session_state.templates.update(templates)
        return True
    except:
        return False

def get_preset_templates():
    return {
        "pitch_deck": {
            "name": "🚀 Startup Pitch Deck",
            "category": "Pitch",
            "slide_count": 10,
            "tone": "Persuasive",
            "audience": "Investors",
            "theme": "Gradient Modern",
            "image_mode": "With Images",
            "language": "English"
        },
        "corporate_report": {
            "name": "📈 Corporate Report",
            "category": "Business",
            "slide_count": 12,
            "tone": "Formal",
            "audience": "Corporate",
            "theme": "Corporate Blue",
            "image_mode": "With Images",
            "language": "English"
        },
        "training_session": {
            "name": "🎓 Training Session",
            "category": "Training",
            "slide_count": 15,
            "tone": "Educational",
            "audience": "Students",
            "theme": "Pastel Soft",
            "image_mode": "With Images",
            "language": "English"
        }
    }

# ============================================================================
# IMAGE GENERATION FUNCTIONS
# ============================================================================

def get_google_image(query, api_key, cx):
    """Get image using Google Custom Search API"""
    try:
        st.session_state.google_searches_used += 1
        url = "https://www.googleapis.com/customsearch/v1"
        params = {
            'key': api_key, 'cx': cx, 'q': query,
            'searchType': 'image', 'num': 3,
            'imgSize': 'large', 'safe': 'active'
        }
        response = requests.get(url, params=params, timeout=10)
        if response.status_code == 200:
            data = response.json()
            if 'items' in data and len(data['items']) > 0:
                for item in data['items'][:3]:
                    try:
                        image_url = item['link']
                        img_response = requests.get(image_url, timeout=10)
                        if img_response.status_code == 200 and len(img_response.content) > 5000:
                            img = Image.open(io.BytesIO(img_response.content))
                            if img.size[0] > 300 and img.size[1] > 200:
                                return img_response.content
                    except:
                        continue
        return None
    except:
        return None

def get_unsplash_image(query, width=800, height=600):
    """Get image from Unsplash"""
    try:
        clean_query = query.strip().replace(' ', ',')
        url = f"https://source.unsplash.com/{width}x{height}/?{clean_query}"
        response = requests.get(url, timeout=15, allow_redirects=True)
        if response.status_code == 200 and len(response.content) > 5000:
            return response.content
        return None
    except:
        return None

def get_topic_relevant_image(main_topic, slide_title, google_api_key, google_cx, use_unsplash):
    """Get relevant image"""
    search_terms = []
    if slide_title:
        search_terms.append(slide_title)
    if main_topic:
        search_terms.append(main_topic)
    
    for term in search_terms:
        if google_api_key and google_cx:
            image_data = get_google_image(term, google_api_key, google_cx)
            if image_data:
                return image_data
        if use_unsplash:
            image_data = get_unsplash_image(term)
            if image_data:
                return image_data
    return None

# ============================================================================
# AI CONTENT GENERATION
# ============================================================================

def repair_truncated_json(json_text):
    """Repair truncated JSON"""
    text = json_text.strip()
    if text.startswith("```json"):
        text = text[7:]
    if text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()
    
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    
    slides = []
    slides_start = text.find('"slides"')
    if slides_start == -1:
        return None
    
    bracket_pos = text.find('[', slides_start)
    if bracket_pos == -1:
        return None
    
    current_pos = bracket_pos + 1
    brace_count = 0
    slide_start = -1
    
    while current_pos < len(text):
        char = text[current_pos]
        if char == '{' and brace_count == 0:
            slide_start = current_pos
            brace_count = 1
        elif char == '{':
            brace_count += 1
        elif char == '}':
            brace_count -= 1
            if brace_count == 0 and slide_start != -1:
                slide_text = text[slide_start:current_pos + 1]
                try:
                    slide_obj = json.loads(slide_text)
                    if 'title' in slide_obj:
                        if 'bullets' not in slide_obj:
                            slide_obj['bullets'] = []
                        if 'image_prompt' not in slide_obj:
                            slide_obj['image_prompt'] = slide_obj['title']
                        if 'speaker_notes' not in slide_obj:
                            slide_obj['speaker_notes'] = ""
                        slides.append(slide_obj)
                except:
                    pass
                slide_start = -1
        current_pos += 1
    
    if slides:
        return {"slides": slides}
    return None

def generate_content_with_ai(api_key, topic, category, slide_count, tone, audience, key_points, model_choice, language, groq_api_key=None):
    """Generate presentation content using AI"""
    try:
        use_groq_api = "Groq" in model_choice and groq_api_key
        
        if use_groq_api:
            model = "llama-3.3-70b-versatile" if "Llama 3.3" in model_choice else "mixtral-8x7b-32768"
            api_url = "https://api.groq.com/openai/v1/chat/completions"
            headers = {"Authorization": f"Bearer {groq_api_key.strip()}", "Content-Type": "application/json"}
        else:
            if "Gemini" in model_choice:
                model = "google/gemini-2.0-flash-exp:free"
            elif "Llama" in model_choice:
                model = "meta-llama/llama-3.2-3b-instruct:free"
            else:
                model = "anthropic/claude-3.5-sonnet"
            api_url = "https://openrouter.ai/api/v1/chat/completions"
            headers = {"Authorization": f"Bearer {api_key.strip()}", "Content-Type": "application/json"}
        
        calculated_tokens = min(slide_count * 350 + 500, 4000)
        language_instruction = f"Generate ALL content in {language} language." if language != "English" else ""
        
        prompt = f"""{language_instruction}
Create a {slide_count}-slide presentation about: {topic}
Category: {category} | Tone: {tone} | Audience: {audience}
{f"Include: {key_points}" if key_points else ""}

Return ONLY JSON:
{{"slides": [
  {{"title": "Title", "bullets": [], "image_prompt": "topic", "speaker_notes": "notes"}},
  {{"title": "Point", "bullets": ["detail 1", "detail 2"], "image_prompt": "topic", "speaker_notes": "notes"}}
]}}

REQUIREMENTS:
1. First slide: TITLE ONLY (empty bullets)
2. Other slides: 3-5 bullets each
3. Total: exactly {slide_count} slides
4. Return ONLY JSON

Generate now:"""

        response = requests.post(api_url, headers=headers,
                                json={"model": model, "max_tokens": calculated_tokens,
                                      "messages": [{"role": "user", "content": prompt}]},
                                timeout=60)
        
        if response.status_code == 200:
            data = response.json()
            content_text = data["choices"][0]["message"]["content"]
            slides_data = repair_truncated_json(content_text)
            
            if slides_data and "slides" in slides_data:
                slides = slides_data["slides"]
                if not slides:
                    return None
                
                for i, slide in enumerate(slides):
                    if 'bullets' not in slide:
                        slide['bullets'] = []
                    if 'image_prompt' not in slide:
                        slide['image_prompt'] = slide.get('title', topic)
                    if 'speaker_notes' not in slide:
                        slide['speaker_notes'] = ""
                
                return slides
        return None
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

# ============================================================================
# POWERPOINT CREATION
# ============================================================================

def create_powerpoint(slides_content, theme, image_mode, google_api_key, google_cx, use_unsplash, category, audience, topic, logo_data, show_progress=True):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    themes = {
        "Corporate Blue": {"bg": RGBColor(240, 248, 255), "accent": RGBColor(31, 119, 180), "text": RGBColor(0, 0, 0)},
        "Gradient Modern": {"bg": RGBColor(240, 242, 246), "accent": RGBColor(138, 43, 226), "text": RGBColor(0, 0, 0)},
        "Minimal Dark": {"bg": RGBColor(30, 30, 30), "accent": RGBColor(255, 215, 0), "text": RGBColor(255, 255, 255)},
        "Pastel Soft": {"bg": RGBColor(255, 250, 240), "accent": RGBColor(255, 182, 193), "text": RGBColor(60, 60, 60)},
        "Professional Green": {"bg": RGBColor(245, 255, 250), "accent": RGBColor(34, 139, 34), "text": RGBColor(0, 0, 0)}
    }
    
    color_scheme = themes.get(theme, themes["Corporate Blue"])
    
    if show_progress:
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    for idx, slide_data in enumerate(slides_content):
        if show_progress:
            status_text.text(f"Creating slide {idx + 1}/{len(slides_content)}...")
            progress_bar.progress((idx + 1) / len(slides_content))
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_scheme["bg"]
        
        if logo_data:
            try:
                logo_stream = io.BytesIO(logo_data)
                slide.shapes.add_picture(logo_stream, Inches(9), Inches(0.2), width=Inches(0.8))
            except:
                pass
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_data["title"]
        title_frame.paragraphs[0].font.size = Pt(36 if idx == 0 else 28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = color_scheme["accent"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.LEFT
        
        if idx > 0 and slide_data.get("bullets"):
            bullet_width = Inches(5.5) if image_mode == "With Images" else Inches(9)
            bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), bullet_width, Inches(4.5))
            text_frame = bullet_box.text_frame
            text_frame.word_wrap = True
            
            for bullet in slide_data["bullets"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = color_scheme["text"]
                p.space_after = Pt(12)
        
        if idx == 0:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = f"{category} Presentation | {audience}"
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].font.color.rgb = color_scheme["text"]
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
        
        if idx > 0 and image_mode == "With Images":
            image_data = get_topic_relevant_image(topic, slide_data["title"], google_api_key, google_cx, use_unsplash)
            if image_data:
                try:
                    image_stream = io.BytesIO(image_data)
                    slide.shapes.add_picture(image_stream, Inches(6.5), Inches(2), width=Inches(3))
                except:
                    pass
            time.sleep(0.3)
    
    if show_progress:
        progress_bar.progress(1.0)
        status_text.text("✅ Presentation created!")
    
    return prs

# ============================================================================
# LOGIN PAGE
# ============================================================================

def show_login_page():
    """Display login page"""
    st.markdown("""
        <style>
        .login-container {
            max-width: 400px;
            margin: 80px auto;
            padding: 40px;
            background: white;
            border-radius: 15px;
            box-shadow: 0 8px 16px rgba(0,0,0,0.1);
        }
        .login-header {
            text-align: center;
            color: #1f77b4;
            margin-bottom: 30px;
            font-size: 28px;
            font-weight: bold;
        }
        </style>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        st.markdown('<div class="login-container">', unsafe_allow_html=True)
        st.markdown('<div style="font-size: 70px; text-align: center; margin-bottom: 20px;">📊</div>', unsafe_allow_html=True)
        st.markdown('<div class="login-header">AI PowerPoint Generator Pro</div>', unsafe_allow_html=True)
        
        st.markdown("### 🔐 Sign In")
        
        username = st.text_input("Username", key="login_username", placeholder="Enter your username")
        password = st.text_input("Password", type="password", key="login_password", placeholder="Enter your password")
        
        col_btn1, col_btn2 = st.columns(2)
        
        with col_btn1:
            if st.button("🔓 Login", use_container_width=True, type="primary"):
                if username and password:
                    user = verify_user(username, password)
                    if user:
                        st.session_state.logged_in = True
                        st.session_state.user = user
                        log_usage(user['id'], 'login')
                        st.success(f"✅ Welcome, {user['full_name']}!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("❌ Invalid username or password")
                else:
                    st.warning("⚠️ Please enter both username and password")
        
        with col_btn2:
            if st.button("🔑 Demo Credentials", use_container_width=True):
                st.info("""
**Admin Login:**
Username: `admin`
Password: `admin123`
                """)
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        st.markdown("---")
        st.markdown("""
            <div style='text-align: center; color: #888; font-size: 12px; margin-top: 20px;'>
                <p>🔒 Secure Login System | © 2025 AI PPT Generator Pro</p>
            </div>
        """, unsafe_allow_html=True)

# ============================================================================
# PAGE CONFIGURATION & INITIALIZATION
# ============================================================================

st.set_page_config(
    page_title="AI PowerPoint Generator Pro",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize database
init_database()

# Initialize session state
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'user' not in st.session_state:
    st.session_state.user = None
if 'generation_count' not in st.session_state:
    st.session_state.generation_count = 0
if 'total_slides' not in st.session_state:
    st.session_state.total_slides = 0
if 'google_searches_used' not in st.session_state:
    st.session_state.google_searches_used = 0
if 'templates' not in st.session_state:
    st.session_state.templates = {}
if 'selected_template' not in st.session_state:
    st.session_state.selected_template = None

# Check login
if not st.session_state.logged_in:
    show_login_page()
    st.stop()

# ============================================================================
# PROFESSIONAL CSS STYLING
# ============================================================================

st.markdown("""
<style>
    .main {
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
    }
    .main-header {
        font-size: 2.8rem;
        font-weight: 800;
        background: linear-gradient(120deg, #1f77b4, #667eea, #764ba2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        text-align: center;
        color: #666;
        font-size: 1.1rem;
        margin-bottom: 2rem;
    }
    .dashboard-card {
        background: white;
        padding: 1.5rem;
        border-radius: 15px;
        box-shadow: 0 4px 20px rgba(0,0,0,0.08);
        transition: transform 0.3s ease;
    }
    .dashboard-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 30px rgba(0,0,0,0.12);
    }
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.2rem;
        border-radius: 12px;
        color: white;
        text-align: center;
    }
    .metric-value {
        font-size: 2rem;
        font-weight: 700;
    }
    .user-info {
        background: #f0f8ff;
        padding: 15px;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    .online-indicator {
        display: inline-block;
        width: 10px;
        height: 10px;
        background: #4caf50;
        border-radius: 50%;
        margin-right: 5px;
        animation: pulse 2s infinite;
    }
    @keyframes pulse {
        0%, 100% { opacity: 1; }
        50% { opacity: 0.5; }
    }
    .activity-card {
        background: #f5f5f5;
        padding: 12px;
        border-radius: 6px;
        margin: 8px 0;
        border-left: 4px solid #1f77b4;
    }
    .form-section {
        background: white;
        padding: 1.5rem;
        border-radius: 15px;
        margin: 1rem 0;
        border-left: 5px solid #667eea;
    }
    .download-section {
        background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%);
        padding: 2rem;
        border-radius: 20px;
        margin: 2rem 0;
        text-align: center;
        color: white;
    }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# SIDEBAR (Common for both Admin & User)
# ============================================================================

with st.sidebar:
    user_stats = get_user_stats(st.session_state.user['id'])
    st.markdown(f"""
    <div class='user-info'>
        <h3>👤 {st.session_state.user['full_name']}</h3>
        <p><b>@{st.session_state.user['username']}</b></p>
        <p>Role: <span style='background:#1f77b4;color:white;padding:2px 8px;border-radius:4px;font-size:12px;'>{st.session_state.user['role'].upper()}</span></p>
        <hr>
        <p>📊 Presentations: <b>{user_stats['total_presentations']}</b></p>
        <p>📄 Total Slides: <b>{user_stats['total_slides']}</b></p>
        <p>🔑 Total Logins: <b>{user_stats['total_logins']}</b></p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.button("🚪 Logout", use_container_width=True, type="primary"):
        logout_user(st.session_state.user['id'])
        log_usage(st.session_state.user['id'], 'logout')
        st.session_state.logged_in = False
        st.session_state.user = None
        st.rerun()
    
    st.markdown("---")
    
    # API Configuration
    st.markdown("### ⚙️ API Configuration")
    
    with st.expander("🔑 AI Models", expanded=True):
        claude_api_key = st.text_input("OpenRouter API Key", type="password", help="For AI content generation")
        
        model_choice = st.selectbox(
            "Select AI Model",
            [
                "Free Model (Google Gemini Flash)",
                "Free Model (Meta Llama 3.2)",
                "Groq (Llama 3.3 70B) - FREE & FAST",
                "Groq (Mixtral 8x7B) - FREE",
                "Claude 3.5 Sonnet (Paid)"
            ]
        )
        
        groq_api_key = None
        if "Groq" in model_choice:
            groq_api_key = st.text_input("Groq API Key (FREE)", type="password", key="groq_key")
            if groq_api_key:
                st.success("✅ Groq configured!")
            st.info("Get FREE key from: https://console.groq.com/")
    
    with st.expander("🖼️ Image Settings"):
        google_api_key = st.text_input("Google API Key", type="password")
        google_cx = st.text_input("Google CX ID")
        if google_api_key and google_cx:
            st.success("✅ Google Search configured!")
        use_unsplash = st.checkbox("Use Unsplash (Free)", value=True)
    
    with st.expander("🏢 Branding"):
        logo_file = st.file_uploader("Company Logo", type=["png", "jpg"])
        logo_data = None
        if logo_file:
            logo_data = logo_file.read()
            st.success("✅ Logo uploaded!")
    
    st.markdown("---")
    
    # Stats Dashboard
    st.markdown("### 📊 Your Statistics")
    col1, col2 = st.columns(2)
    with col1:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-value">{user_stats['total_presentations']}</div>
            <div style='font-size:0.85rem;opacity:0.9;'>Presentations</div>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-value">{user_stats['total_slides']}</div>
            <div style='font-size:0.85rem;opacity:0.9;'>Total Slides</div>
        </div>
        """, unsafe_allow_html=True)

# ============================================================================
# ADMIN DASHBOARD
# ============================================================================

if st.session_state.user['role'] == 'admin':
    
    st.markdown('<div class="main-header">👑 Admin Dashboard</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Complete System Management & PPT Generator</div>', unsafe_allow_html=True)
    
    # System Statistics
    sys_stats = get_system_stats()
    
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.metric("👥 Total Users", sys_stats['total_users'])
    with col2:
        st.metric("🟢 Online Now", sys_stats['currently_online'])
    with col3:
        st.metric("📊 Presentations", sys_stats['total_presentations'])
    with col4:
        st.metric("📄 Total Slides", sys_stats['total_slides'])
    with col5:
        st.metric("🕒 Today Logins", sys_stats['today_logins'])
    
    st.markdown("---")
    
    # Main Navigation Tabs
    main_tab = st.radio(
        "Navigation",
        ["📊 PPT Generator Dashboard", "👥 User Management", "📈 System Analytics"],
        horizontal=True
    )
    
    # ====================
    # TAB 1: PPT GENERATOR DASHBOARD (Full Features)
    # ====================
    
    if main_tab == "📊 PPT Generator Dashboard":
        st.markdown("## 🚀 AI PowerPoint Generator")
        
        ppt_tabs = st.tabs(["📝 Create", "📁 Templates", "📜 History", "⚙️ Settings"])
        
        # CREATE TAB
        with ppt_tabs[0]:
            st.subheader("🎯 Create New Presentation")
            
            # Quick Templates
            st.markdown("### 🚀 Quick Start Templates")
            preset_templates = get_preset_templates()
            cols = st.columns(3)
            
            for idx, (key, template) in enumerate(preset_templates.items()):
                with cols[idx % 3]:
                    if st.button(f"{template['name']}", key=f"admin_preset_{key}", use_container_width=True):
                        st.session_state.selected_template = template
            
            st.markdown("---")
            
            # Main Form
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown('<div class="form-section">', unsafe_allow_html=True)
                st.markdown("📝 **Content Details**", unsafe_allow_html=True)
                
                topic = st.text_input("Topic *", placeholder="e.g., AI in Healthcare", key="admin_topic")
                
                categories = ["Business", "Pitch", "Marketing", "Technical", "Academic", "Training"]
                category = st.selectbox("Category *", categories, key="admin_category")
                
                col1_1, col1_2 = st.columns(2)
                with col1_1:
                    slide_count = st.number_input("Slides *", min_value=3, max_value=20, value=6, key="admin_slides")
                with col1_2:
                    languages = ["English", "Hindi", "Spanish", "French"]
                    language = st.selectbox("Language", languages, key="admin_lang")
                
                tones = ["Formal", "Neutral", "Inspirational", "Educational", "Persuasive"]
                tone = st.selectbox("Tone *", tones, key="admin_tone")
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            with col2:
                st.markdown('<div class="form-section">', unsafe_allow_html=True)
                st.markdown("🎨 **Design & Style**", unsafe_allow_html=True)
                
                audiences = ["Investors", "Students", "Corporate", "Clients", "Managers"]
                audience = st.selectbox("Target Audience *", audiences, key="admin_audience")
                
                themes_list = ["Corporate Blue", "Gradient Modern", "Minimal Dark", "Pastel Soft", "Professional Green"]
                theme = st.selectbox("Visual Theme *", themes_list, key="admin_theme")
                
                image_modes = ["With Images", "No Images"]
                image_mode = st.selectbox("Image Mode *", image_modes, key="admin_imgmode")
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            with st.expander("➕ Additional Options"):
                key_points = st.text_area("Key Points", placeholder="- Point 1\n- Point 2", key="admin_keypoints")
            
            st.markdown("---")
            
            if st.button("🚀 Generate Presentation", use_container_width=True, type="primary", key="admin_generate"):
                if topic:
                    has_valid_api = False
                    if "Groq" in model_choice and groq_api_key:
                        has_valid_api = True
                    elif claude_api_key:
                        has_valid_api = True
                    
                    if has_valid_api:
                        with st.spinner("🤖 Generating your presentation..."):
                            slides_content = generate_content_with_ai(
                                claude_api_key, topic, category, slide_count,
                                tone, audience, key_points, model_choice, language,
                                groq_api_key=groq_api_key
                            )
                            
                            if slides_content:
                                log_usage(st.session_state.user['id'], 'generate_presentation', topic, len(slides_content))
                                
                                prs = create_powerpoint(
                                    slides_content, theme, image_mode,
                                    google_api_key if 'google_api_key' in locals() else "",
                                    google_cx if 'google_cx' in locals() else "",
                                    use_unsplash, category, audience, topic, logo_data
                                )
                                
                                pptx_io = io.BytesIO()
                                prs.save(pptx_io)
                                pptx_io.seek(0)
                                
                                st.markdown("""
                                <div class="download-section">
                                    <h2>🎉 Presentation Ready!</h2>
                                    <p>Download your professionally generated presentation</p>
                                </div>
                                """, unsafe_allow_html=True)
                                
                                st.download_button(
                                    "📥 DOWNLOAD POWERPOINT",
                                    pptx_io.getvalue(),
                                    f"{topic.replace(' ', '_')}.pptx",
                                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True,
                                    type="primary"
                                )
                                
                                st.balloons()
                                
                                with st.expander("📄 Preview Slides", expanded=True):
                                    for idx, slide in enumerate(slides_content):
                                        st.markdown(f"### Slide {idx + 1}: {slide['title']}")
                                        if slide.get('bullets'):
                                            for bullet in slide['bullets']:
                                                st.markdown(f"• {bullet}")
                                        st.markdown("---")
                    else:
                        st.error("⚠️ Please configure API keys in sidebar")
                else:
                    st.error("⚠️ Please enter a topic")
        
        # TEMPLATES TAB
        with ppt_tabs[1]:
            st.subheader("📁 Template Manager")
            if st.session_state.templates:
                for temp_id, template in st.session_state.templates.items():
                    st.markdown(f"**{template['name']}** - {template.get('category', 'General')}")
                    if st.button("Use Template", key=f"admin_use_{temp_id}"):
                        st.session_state.selected_template = template
            else:
                st.info("No templates saved yet. Create one from the 'Create' tab!")
        
        # HISTORY TAB
        with ppt_tabs[2]:
            st.subheader("📜 Your Generation History")
            my_activities = get_user_activity_details(st.session_state.user['id'])
            if my_activities:
                for activity in my_activities[:10]:
                    action, topic, slides_count, timestamp = activity
                    if action == 'generate_presentation':
                        st.markdown(f"""
<div class='activity-card'>
    <b>📊 {topic}</b><br>
    <small>📄 {slides_count} slides | 🕒 {timestamp}</small>
</div>
                        """, unsafe_allow_html=True)
            else:
                st.info("No history yet")
        
        # SETTINGS TAB
        with ppt_tabs[3]:
            st.subheader("⚙️ Generator Settings")
            st.info("Configure your default preferences here")
    
    # ====================
    # TAB 2: USER MANAGEMENT (Create/Edit/Delete)
    # ====================
    
    elif main_tab == "👥 User Management":
        st.markdown("## 👥 User Management System")
        
        user_mgmt_tabs = st.tabs([
            "🟢 Active Users",
            "➕ Create User",
            "📋 All Users",
            "✏️ Edit User",
            "🗑️ Delete User"
        ])
        
        # ACTIVE USERS TAB
        with user_mgmt_tabs[0]:
            st.subheader("🟢 Currently Logged In Users")
            
            if st.button("🔄 Refresh", key="refresh_active"):
                st.rerun()
            
            active_users = get_currently_logged_in_users()
            
            if active_users:
                st.success(f"**{len(active_users)} user(s) currently online**")
                
                for user in active_users:
                    user_id, username, email, full_name, login_time, role = user
                    user_stats_temp = get_user_stats(user_id)
                    
                    st.markdown(f"""
<div style='background:#e8f5e9;padding:15px;border-radius:10px;margin:10px 0;'>
    <span class='online-indicator'></span>
    <b style='font-size:18px;'>{full_name}</b> <span style='background:#4caf50;color:white;padding:2px 8px;border-radius:4px;font-size:12px;'>{role}</span><br>
    <small>👤 @{username} | 📧 {email if email else 'N/A'}</small><br>
    <small>🕒 Logged in: {login_time}</small><br>
    <small>📊 {user_stats_temp['total_presentations']} presentations | 📄 {user_stats_temp['total_slides']} slides</small>
</div>
                    """, unsafe_allow_html=True)
            else:
                st.warning("⚠️ No users currently online")
        
        # CREATE USER TAB
        with user_mgmt_tabs[1]:
            st.subheader("➕ Create New User")
            
            with st.form("create_user_form"):
                col_a, col_b = st.columns(2)
                
                with col_a:
                    new_username = st.text_input("Username *", placeholder="john_doe")
                    new_full_name = st.text_input("Full Name *", placeholder="John Doe")
                    new_email = st.text_input("Email", placeholder="john@company.com")
                
                with col_b:
                    new_phone = st.text_input("Phone", placeholder="+1 234 567 8900")
                    new_password = st.text_input("Password *", type="password", placeholder="Min 6 characters")
                    confirm_password = st.text_input("Confirm Password *", type="password")
                
                new_role = st.selectbox("Role", ["user", "admin"])
                
                submitted = st.form_submit_button("✅ Create User", use_container_width=True, type="primary")
                
                if submitted:
                    if new_username and new_password and new_full_name:
                        if new_password == confirm_password:
                            if len(new_password) >= 6:
                                success, message = create_user(new_username, new_password, new_email, new_full_name, new_phone, new_role)
                                if success:
                                    st.success(f"""
### ✅ User Created Successfully!

**Credentials to share:**
- Username: `{new_username}`
- Password: `{new_password}`
- Full Name: {new_full_name}
- Role: {new_role}

📧 Send these credentials to the user securely.
                                    """)
                                    time.sleep(2)
                                    st.rerun()
                                else:
                                    st.error(f"❌ {message}")
                            else:
                                st.error("❌ Password must be at least 6 characters")
                        else:
                            st.error("❌ Passwords don't match")
                    else:
                        st.warning("⚠️ Please fill all required fields (marked with *)")
        
        # ALL USERS TAB
        with user_mgmt_tabs[2]:
            st.subheader("📋 All Registered Users")
            
            users = get_all_users()
            
            user_data = []
            for user in users:
                user_data.append({
                    'ID': user[0],
                    'Username': user[1],
                    'Full Name': user[3] if user[3] else user[1],
                    'Email': user[2] if user[2] else 'N/A',
                    'Phone': user[4] if user[4] else 'N/A',
                    'Created': user[5],
                    'Last Login': user[6] if user[6] else 'Never',
                    'Status': '✅ Active' if user[7] else '❌ Inactive',
                    'Role': user[8]
                })
            
            df_users = pd.DataFrame(user_data)
            st.dataframe(df_users, use_container_width=True, height=500)
            
            # Download as CSV
            csv = df_users.to_csv(index=False)
            st.download_button(
                "📥 Download User List (CSV)",
                csv,
                f"users_{datetime.now().strftime('%Y%m%d')}.csv",
                "text/csv",
                key='download-users'
            )
        
        # EDIT USER TAB
        with user_mgmt_tabs[3]:
            st.subheader("✏️ Edit User")
            
            all_users_list = get_all_users()
            user_options = {f"{u[0]} - {u[1]} ({u[3] if u[3] else u[1]})": u[0] for u in all_users_list}
            
            selected_user_display = st.selectbox("Select User to Edit", list(user_options.keys()))
            
            if selected_user_display:
                selected_user_id = user_options[selected_user_display]
                
                if selected_user_id == 1:
                    st.warning("⚠️ Cannot edit the default admin account for security reasons")
                else:
                    user_details = get_user_by_id(selected_user_id)
                    
                    if user_details:
                        with st.form("edit_user_form"):
                            st.info(f"Editing User ID: {user_details[0]}")
                            
                            col_e1, col_e2 = st.columns(2)
                            
                            with col_e1:
                                edit_username = st.text_input("Username *", value=user_details[1])
                                edit_full_name = st.text_input("Full Name *", value=user_details[3] if user_details[3] else "")
                                edit_email = st.text_input("Email", value=user_details[2] if user_details[2] else "")
                            
                            with col_e2:
                                edit_phone = st.text_input("Phone", value=user_details[4] if user_details[4] else "")
                                edit_is_active = st.selectbox("Status", [1, 0], format_func=lambda x: "✅ Active" if x == 1 else "❌ Inactive", index=0 if user_details[5] else 1)
                                edit_password = st.text_input("New Password (leave blank to keep current)", type="password")
                            
                            update_submitted = st.form_submit_button("💾 Update User", use_container_width=True, type="primary")
                            
                            if update_submitted:
                                if edit_username and edit_full_name:
                                    success, message = update_user(
                                        selected_user_id,
                                        edit_username,
                                        edit_email,
                                        edit_full_name,
                                        edit_phone,
                                        edit_is_active,
                                        edit_password if edit_password else None
                                    )
                                    if success:
                                        st.success(f"✅ {message}")
                                        time.sleep(2)
                                        st.rerun()
                                    else:
                                        st.error(f"❌ {message}")
                                else:
                                    st.warning("⚠️ Username and Full Name are required")
        
        # DELETE USER TAB
        with user_mgmt_tabs[4]:
            st.subheader("🗑️ Delete User")
            
            st.warning("⚠️ **Warning:** Deleting a user is permanent and cannot be undone!")
            
            all_users_del = get_all_users()
            user_del_options = {f"{u[0]} - {u[1]} ({u[3] if u[3] else u[1]})": u[0] for u in all_users_del}
            
            selected_del_user = st.selectbox("Select User to Delete", list(user_del_options.keys()), key="del_user_select")
            
            if selected_del_user:
                del_user_id = user_del_options[selected_del_user]
                
                if del_user_id == 1:
                    st.error("❌ Cannot delete the default admin account!")
                else:
                    user_del_details = get_user_by_id(del_user_id)
                    
                    if user_del_details:
                        st.markdown(f"""
### User Details:
- **ID:** {user_del_details[0]}
- **Username:** {user_del_details[1]}
- **Full Name:** {user_del_details[3]}
- **Email:** {user_del_details[2]}
- **Role:** {user_del_details[6]}
                        """)
                        
                        confirm_delete = st.checkbox("I confirm I want to delete this user")
                        
                        if confirm_delete:
                            if st.button("🗑️ PERMANENTLY DELETE USER", type="primary", use_container_width=True):
                                success, message = delete_user(del_user_id)
                                if success:
                                    st.success(f"✅ {message}")
                                    time.sleep(2)
                                    st.rerun()
                                else:
                                    st.error(f"❌ {message}")
    
    # ====================
    # TAB 3: SYSTEM ANALYTICS
    # ====================
    
    elif main_tab == "📈 System Analytics":
        st.markdown("## 📈 System Analytics Dashboard")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("📊 User Growth Over Time")
            # Placeholder for user growth chart
            st.info("User growth analytics coming soon")
        
        with col2:
            st.subheader("📄 Presentation Generation Trends")
            # Placeholder for generation trends
            st.info("Generation trends analytics coming soon")
        
        st.markdown("---")
        
        # Activity Log
        st.subheader("📊 Recent System Activities (Last 50)")
        all_activities = get_all_user_activities()
        
        if all_activities:
            activity_records = []
            for activity in all_activities[:50]:
                username, action, topic, slides_count, timestamp = activity
                activity_records.append({
                    'Username': username,
                    'Action': action,
                    'Topic': topic if topic else '-',
                    'Slides': slides_count if slides_count else '-',
                    'Timestamp': timestamp
                })
            
            df_activities = pd.DataFrame(activity_records)
            st.dataframe(df_activities, use_container_width=True, height=600)
            
            # Download activity log
            csv_act = df_activities.to_csv(index=False)
            st.download_button(
                "📥 Download Activity Log (CSV)",
                csv_act,
                f"activity_log_{datetime.now().strftime('%Y%m%d')}.csv",
                "text/csv"
            )

# ============================================================================
# USER DASHBOARD
# ============================================================================

else:
    
    st.markdown('<div class="main-header">📊 AI PowerPoint Generator Pro</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Create stunning presentations with AI-powered content</div>', unsafe_allow_html=True)
    
    # User Tabs
    user_tabs = st.tabs(["📝 Create Presentation", "📁 My Templates", "📜 My History", "⚙️ Settings"])
    
    # CREATE TAB
    with user_tabs[0]:
        st.subheader("🎯 Create New Presentation")
        
        # Quick Templates
        st.markdown("### 🚀 Quick Start Templates")
        preset_templates = get_preset_templates()
        cols = st.columns(3)
        
        for idx, (key, template) in enumerate(preset_templates.items()):
            with cols[idx % 3]:
                if st.button(f"{template['name']}", key=f"user_preset_{key}", use_container_width=True):
                    st.session_state.selected_template = template
        
        st.markdown("---")
        
        # Main Form
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown('<div class="form-section">', unsafe_allow_html=True)
            st.markdown("📝 **Content Details**", unsafe_allow_html=True)
            
            topic = st.text_input("Topic *", placeholder="e.g., Digital Marketing Strategy", key="user_topic")
            
            categories = ["Business", "Pitch", "Marketing", "Technical", "Academic", "Training"]
            category = st.selectbox("Category *", categories, key="user_category")
            
            col1_1, col1_2 = st.columns(2)
            with col1_1:
                slide_count = st.number_input("Slides *", min_value=3, max_value=20, value=6, key="user_slides")
            with col1_2:
                languages = ["English", "Hindi", "Spanish", "French"]
                language = st.selectbox("Language", languages, key="user_lang")
            
            tones = ["Formal", "Neutral", "Inspirational", "Educational", "Persuasive"]
            tone = st.selectbox("Tone *", tones, key="user_tone")
            
            st.markdown('</div>', unsafe_allow_html=True)
        
        with col2:
            st.markdown('<div class="form-section">', unsafe_allow_html=True)
            st.markdown("🎨 **Design & Style**", unsafe_allow_html=True)
            
            audiences = ["Investors", "Students", "Corporate", "Clients", "Managers"]
            audience = st.selectbox("Target Audience *", audiences, key="user_audience")
            
            themes_list = ["Corporate Blue", "Gradient Modern", "Minimal Dark", "Pastel Soft", "Professional Green"]
            theme = st.selectbox("Visual Theme *", themes_list, key="user_theme")
            
            image_modes = ["With Images", "No Images"]
            image_mode = st.selectbox("Image Mode *", image_modes, key="user_imgmode")
            
            st.markdown('</div>', unsafe_allow_html=True)
        
        with st.expander("➕ Additional Options"):
            key_points = st.text_area("Key Points", placeholder="- Important point 1\n- Important point 2", key="user_keypoints")
        
        st.markdown("---")
        
        if st.button("🚀 Generate Presentation", use_container_width=True, type="primary", key="user_generate"):
            if topic:
                has_valid_api = False
                if "Groq" in model_choice and groq_api_key:
                    has_valid_api = True
                elif claude_api_key:
                    has_valid_api = True
                
                if has_valid_api:
                    with st.spinner("🤖 Generating your presentation..."):
                        slides_content = generate_content_with_ai(
                            claude_api_key, topic, category, slide_count,
                            tone, audience, key_points, model_choice, language,
                            groq_api_key=groq_api_key
                        )
                        
                        if slides_content:
                            log_usage(st.session_state.user['id'], 'generate_presentation', topic, len(slides_content))
                            
                            prs = create_powerpoint(
                                slides_content, theme, image_mode,
                                google_api_key if 'google_api_key' in locals() else "",
                                google_cx if 'google_cx' in locals() else "",
                                use_unsplash, category, audience, topic, logo_data
                            )
                            
                            pptx_io = io.BytesIO()
                            prs.save(pptx_io)
                            pptx_io.seek(0)
                            
                            st.markdown("""
                            <div class="download-section">
                                <h2>🎉 Your Presentation is Ready!</h2>
                                <p>Download your professionally generated presentation below</p>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            st.download_button(
                                "📥 DOWNLOAD POWERPOINT",
                                pptx_io.getvalue(),
                                f"{topic.replace(' ', '_')}.pptx",
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                                type="primary"
                            )
                            
                            st.balloons()
                            
                            with st.expander("📄 Preview Generated Content", expanded=True):
                                for idx, slide in enumerate(slides_content):
                                    st.markdown(f"### Slide {idx + 1}: {slide['title']}")
                                    if slide.get('bullets'):
                                        for bullet in slide['bullets']:
                                            st.markdown(f"• {bullet}")
                                    st.markdown("---")
                else:
                    st.error("⚠️ Please configure API keys in sidebar to generate presentations")
            else:
                st.error("⚠️ Please enter a presentation topic")
    
    # TEMPLATES TAB
    with user_tabs[1]:
        st.subheader("📁 My Saved Templates")
        if st.session_state.templates:
            for temp_id, template in st.session_state.templates.items():
                col_t1, col_t2 = st.columns([3, 1])
                with col_t1:
                    st.markdown(f"**{template['name']}** - {template.get('category', 'General')}")
                with col_t2:
                    if st.button("Use", key=f"user_use_{temp_id}"):
                        st.session_state.selected_template = template
                        st.rerun()
        else:
            st.info("💡 No templates saved yet. Create presentations and save them as templates!")
    
    # HISTORY TAB
    with user_tabs[2]:
        st.subheader("📜 My Generation History")
        my_activities = get_user_activity_details(st.session_state.user['id'])
        if my_activities:
            for activity in my_activities:
                action, topic, slides_count, timestamp = activity
                if action == 'generate_presentation':
                    st.markdown(f"""
<div class='activity-card'>
    <b>📊 {topic}</b><br>
    <small>📄 {slides_count} slides | 🕒 {timestamp}</small>
</div>
                    """, unsafe_allow_html=True)
        else:
            st.info("📊 No presentations generated yet. Create your first one!")
    
    # SETTINGS TAB
    with user_tabs[3]:
        st.subheader("⚙️ My Settings")
        
        st.markdown("### 👤 Profile Information")
        st.info(f"""
**Full Name:** {st.session_state.user['full_name']}  
**Username:** @{st.session_state.user['username']}  
**Email:** {st.session_state.user['email'] if st.session_state.user['email'] else 'Not set'}  
**Role:** {st.session_state.user['role'].title()}  
        """)
        
        st.markdown("---")
        st.markdown("### 📊 My Statistics")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Total Presentations", user_stats['total_presentations'])
        with col2:
            st.metric("Total Slides", user_stats['total_slides'])
        with col3:
            st.metric("Total Logins", user_stats['total_logins'])

# ============================================================================
# FOOTER
# ============================================================================

st.markdown("---")
st.markdown(f"""
<div style='text-align: center; color: #666; padding: 1rem 0;'>
    <p>🔐 Logged in as: <b>{st.session_state.user['full_name']}</b> (@{st.session_state.user['username']}) | 
    Role: <b>{st.session_state.user['role'].upper()}</b></p>
    <p>✨ AI PowerPoint Generator Pro | Version 3.0 | © 2025</p>
    <p><small>{datetime.now().strftime("%A, %B %d, %Y - %H:%M:%S")}</small></p>
</div>
""", unsafe_allow_html=True)
